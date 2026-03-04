"""
Multi-format export engine for the Multimodal Librarian system.

This module provides comprehensive export functionality supporting multiple formats
including .txt, .docx, .pdf, .rtf, .pptx, .xlsx, and .html with multimedia content preservation.
"""

import io
import logging
import os
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, BinaryIO, Dict, List, Optional, Union

# Document generation libraries
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.shared import OxmlElement, qn
from docx.shared import Inches, Pt
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.drawing.image import Image as XLImage
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PPTXInches
from pptx.util import Pt as PPTXPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import A4, letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image as RLImage
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from ...models.core import (
    ExportMetadata,
    KnowledgeCitation,
    MediaElement,
    MultimediaResponse,
    Visualization,
)

logger = logging.getLogger(__name__)


class ExportEngine:
    """
    Multi-format export engine supporting various document formats.
    
    Supports export to:
    - .txt: Clean text with citations
    - .docx: Word document with embedded media
    - .pdf: PDF with layout and multimedia elements
    - .rtf: Rich Text Format with formatting
    - .pptx: PowerPoint presentation with slides
    - .xlsx: Excel spreadsheet with structured data
    - .html: HTML for web-compatible output
    """
    
    def __init__(self):
        """Initialize the export engine."""
        self.supported_formats = {
            'txt': self._export_to_txt,
            'docx': self._export_to_docx,
            'pdf': self._export_to_pdf,
            'rtf': self._export_to_rtf,
            'pptx': self._export_to_pptx,
            'xlsx': self._export_to_xlsx,
            'html': self._export_to_html
        }
        
    def export_to_format(self, response: MultimediaResponse, format_type: str) -> bytes:
        """
        Export multimedia response to specified format.
        
        Args:
            response: MultimediaResponse to export
            format_type: Target format ('txt', 'docx', 'pdf', 'rtf', 'pptx', 'xlsx')
            
        Returns:
            bytes: Exported document content
            
        Raises:
            ValueError: If format is not supported
            Exception: If export fails
        """
        format_type = format_type.lower().replace('.', '')
        
        if format_type not in self.supported_formats:
            raise ValueError(f"Unsupported format: {format_type}. Supported formats: {list(self.supported_formats.keys())}")
        
        try:
            logger.info(f"Starting export to {format_type} format")
            export_func = self.supported_formats[format_type]
            content = export_func(response)
            
            # Update export metadata
            if response.export_metadata is None:
                response.export_metadata = ExportMetadata(export_format=format_type)
            else:
                response.export_metadata.export_format = format_type
                response.export_metadata.created_at = datetime.now()
            
            response.export_metadata.file_size = len(content)
            response.export_metadata.includes_media = response.has_multimedia()
            
            logger.info(f"Successfully exported to {format_type}, size: {len(content)} bytes")
            return content
            
        except Exception as e:
            logger.error(f"Export to {format_type} failed: {str(e)}")
            raise
    
    def _export_to_txt(self, response: MultimediaResponse) -> bytes:
        """Export to clean text format with citations."""
        logger.debug("Exporting to TXT format")
        
        content_lines = []
        
        # Add main text content
        content_lines.append("MULTIMODAL LIBRARIAN EXPORT")
        content_lines.append("=" * 50)
        content_lines.append(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        content_lines.append("")
        
        # Add main content
        content_lines.append("CONTENT:")
        content_lines.append("-" * 20)
        content_lines.append(response.text_content)
        content_lines.append("")
        
        # Add visualizations descriptions
        if response.visualizations:
            content_lines.append("VISUALIZATIONS:")
            content_lines.append("-" * 20)
            for i, viz in enumerate(response.visualizations, 1):
                content_lines.append(f"{i}. {viz.viz_type.upper()}: {viz.caption}")
                if viz.alt_text:
                    content_lines.append(f"   Description: {viz.alt_text}")
                content_lines.append("")
        
        # Add audio/video descriptions
        if response.audio_content:
            content_lines.append("AUDIO CONTENT:")
            content_lines.append("-" * 20)
            content_lines.append(f"Duration: {response.audio_content.duration_seconds:.1f} seconds")
            content_lines.append(f"Format: {response.audio_content.format}")
            content_lines.append("")
        
        if response.video_content:
            content_lines.append("VIDEO CONTENT:")
            content_lines.append("-" * 20)
            content_lines.append(f"Duration: {response.video_content.duration_seconds:.1f} seconds")
            content_lines.append(f"Resolution: {response.video_content.resolution}")
            content_lines.append(f"Format: {response.video_content.format}")
            content_lines.append("")
        
        # Add citations
        if response.knowledge_citations:
            content_lines.append("CITATIONS:")
            content_lines.append("-" * 20)
            for i, citation in enumerate(response.knowledge_citations, 1):
                content_lines.append(f"{i}. {citation.source_title}")
                content_lines.append(f"   Source: {citation.source_type.value.title()}")
                content_lines.append(f"   Location: {citation.location_reference}")
                content_lines.append(f"   Relevance: {citation.relevance_score:.2f}")
                content_lines.append("")
        
        return "\n".join(content_lines).encode('utf-8')
    
    def _export_to_docx(self, response: MultimediaResponse) -> bytes:
        """Export to Word document with embedded media."""
        logger.debug("Exporting to DOCX format")
        
        doc = Document()
        
        # Add title
        title = doc.add_heading('Multimodal Librarian Export', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add metadata
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        doc.add_paragraph("")
        
        # Add main content
        content_heading = doc.add_heading('Content', level=1)
        content_para = doc.add_paragraph(response.text_content)
        content_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        
        # Add visualizations
        if response.visualizations:
            viz_heading = doc.add_heading('Visualizations', level=1)
            
            for i, viz in enumerate(response.visualizations, 1):
                doc.add_heading(f'{i}. {viz.viz_type.title()}', level=2)
                
                # Add image if we have content data
                if viz.content_data:
                    try:
                        # Save image to temporary file
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                            tmp_file.write(viz.content_data)
                            tmp_file.flush()
                            
                            # Add image to document
                            doc.add_picture(tmp_file.name, width=Inches(6))
                            
                            # Clean up temp file
                            os.unlink(tmp_file.name)
                    except Exception as e:
                        logger.warning(f"Could not embed visualization {i}: {e}")
                        doc.add_paragraph(f"[Visualization: {viz.viz_type}]")
                
                # Add caption
                if viz.caption:
                    caption_para = doc.add_paragraph(viz.caption)
                    caption_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    caption_para.italic = True
                
                doc.add_paragraph("")
        
        # Add multimedia descriptions
        if response.audio_content or response.video_content:
            media_heading = doc.add_heading('Multimedia Content', level=1)
            
            if response.audio_content:
                doc.add_heading('Audio', level=2)
                doc.add_paragraph(f"Duration: {response.audio_content.duration_seconds:.1f} seconds")
                doc.add_paragraph(f"Format: {response.audio_content.format}")
                doc.add_paragraph("")
            
            if response.video_content:
                doc.add_heading('Video', level=2)
                doc.add_paragraph(f"Duration: {response.video_content.duration_seconds:.1f} seconds")
                doc.add_paragraph(f"Resolution: {response.video_content.resolution}")
                doc.add_paragraph(f"Format: {response.video_content.format}")
                doc.add_paragraph("")
        
        # Add citations
        if response.knowledge_citations:
            citations_heading = doc.add_heading('Citations', level=1)
            
            for i, citation in enumerate(response.knowledge_citations, 1):
                citation_para = doc.add_paragraph()
                citation_para.add_run(f"{i}. ").bold = True
                citation_para.add_run(citation.source_title).bold = True
                citation_para.add_run(f"\n   Source: {citation.source_type.value.title()}")
                citation_para.add_run(f"\n   Location: {citation.location_reference}")
                citation_para.add_run(f"\n   Relevance: {citation.relevance_score:.2f}")
                doc.add_paragraph("")
        
        # Save to bytes
        doc_buffer = io.BytesIO()
        doc.save(doc_buffer)
        doc_buffer.seek(0)
        return doc_buffer.getvalue()
    
    def _export_to_pdf(self, response: MultimediaResponse) -> bytes:
        """Export to PDF with layout and multimedia elements."""
        logger.debug("Exporting to PDF format")
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=14,
            spaceAfter=12,
            spaceBefore=20
        )
        
        # Add title
        story.append(Paragraph("Multimodal Librarian Export", title_style))
        story.append(Spacer(1, 12))
        
        # Add metadata
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))
        story.append(Spacer(1, 20))
        
        # Add main content - parse markdown-style formatting
        import re
        def _md_to_html(text):
            """Convert markdown bold and links to ReportLab HTML."""
            # Convert markdown links [text](url) to clickable <a> tags
            text = re.sub(r'\[([^\]]+)\]\(([^)]+)\)', r'<a href="\2" color="blue">[\1]</a>', text)
            # Convert bold **text** to <b>text</b>
            text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
            return text
        
        for line in response.text_content.split('\n'):
            stripped = line.strip()
            if not stripped:
                story.append(Spacer(1, 6))
            elif stripped.startswith('# '):
                story.append(Paragraph(stripped[2:], title_style))
            elif stripped.startswith('## '):
                story.append(Spacer(1, 10))
                story.append(Paragraph(stripped[3:], heading_style))
            else:
                story.append(Paragraph(_md_to_html(stripped), styles['Normal']))
        story.append(Spacer(1, 20))
        
        # Add visualizations
        if response.visualizations:
            story.append(Paragraph("Visualizations", heading_style))
            
            for i, viz in enumerate(response.visualizations, 1):
                story.append(Paragraph(f"{i}. {viz.viz_type.title()}", styles['Heading3']))
                
                # Add image if we have content data
                if viz.content_data:
                    try:
                        # Create image from bytes directly
                        img_buffer = io.BytesIO(viz.content_data)
                        img = RLImage(img_buffer, width=6*inch, height=4*inch)
                        story.append(img)
                    except Exception as e:
                        logger.warning(f"Could not embed visualization {i}: {e}")
                        story.append(Paragraph(f"[Visualization: {viz.viz_type}]", styles['Normal']))
                
                # Add caption
                if viz.caption:
                    caption_style = ParagraphStyle(
                        'Caption',
                        parent=styles['Normal'],
                        fontSize=10,
                        alignment=TA_CENTER,
                        fontName='Helvetica-Oblique'
                    )
                    story.append(Paragraph(viz.caption, caption_style))
                
                story.append(Spacer(1, 12))
        
        # Add multimedia descriptions
        if response.audio_content or response.video_content:
            story.append(Paragraph("Multimedia Content", heading_style))
            
            if response.audio_content:
                story.append(Paragraph("Audio", styles['Heading3']))
                story.append(Paragraph(f"Duration: {response.audio_content.duration_seconds:.1f} seconds", styles['Normal']))
                story.append(Paragraph(f"Format: {response.audio_content.format}", styles['Normal']))
                story.append(Spacer(1, 12))
            
            if response.video_content:
                story.append(Paragraph("Video", styles['Heading3']))
                story.append(Paragraph(f"Duration: {response.video_content.duration_seconds:.1f} seconds", styles['Normal']))
                story.append(Paragraph(f"Resolution: {response.video_content.resolution}", styles['Normal']))
                story.append(Paragraph(f"Format: {response.video_content.format}", styles['Normal']))
                story.append(Spacer(1, 12))
        
        # Add citations
        if response.knowledge_citations:
            story.append(Paragraph("Citations", heading_style))
            
            for i, citation in enumerate(response.knowledge_citations, 1):
                citation_text = f"{i}. <b>{citation.source_title}</b><br/>"
                citation_text += f"Source: {citation.source_type.value.title()}<br/>"
                citation_text += f"Location: {citation.location_reference}<br/>"
                citation_text += f"Relevance: {citation.relevance_score:.2f}"
                
                story.append(Paragraph(citation_text, styles['Normal']))
                story.append(Spacer(1, 8))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()
    
    def _export_to_rtf(self, response: MultimediaResponse) -> bytes:
        """Export to RTF with rich formatting."""
        logger.debug("Exporting to RTF format")
        
        # Create RTF content manually (basic RTF format)
        rtf_content = []
        
        # RTF header
        rtf_content.append(r"{\rtf1\ansi\deff0")
        rtf_content.append(r"{\fonttbl{\f0 Times New Roman;}}")
        rtf_content.append(r"{\colortbl;\red0\green0\blue0;}")
        
        # Title
        rtf_content.append(r"\f0\fs28\b Multimodal Librarian Export\b0\fs24\par")
        rtf_content.append(r"\par")
        
        # Metadata
        rtf_content.append(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\\par")
        rtf_content.append(r"\par")
        
        # Content heading
        rtf_content.append(r"\fs20\b Content\b0\fs18\par")
        
        # Main content - escape RTF special characters
        content_escaped = response.text_content.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}')
        rtf_content.append(f"{content_escaped}\\par")
        rtf_content.append(r"\par")
        
        # Visualizations
        if response.visualizations:
            rtf_content.append(r"\fs20\b Visualizations\b0\fs18\par")
            
            for i, viz in enumerate(response.visualizations, 1):
                rtf_content.append(f"\\b {i}. {viz.viz_type.title()}\\b0\\par")
                
                if viz.caption:
                    caption_escaped = viz.caption.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}')
                    rtf_content.append(f"\\i {caption_escaped}\\i0\\par")
                
                rtf_content.append(r"\par")
        
        # Multimedia content
        if response.audio_content or response.video_content:
            rtf_content.append(r"\fs20\b Multimedia Content\b0\fs18\par")
            
            if response.audio_content:
                rtf_content.append(f"Audio: {response.audio_content.duration_seconds:.1f}s, {response.audio_content.format}\\par")
            
            if response.video_content:
                rtf_content.append(f"Video: {response.video_content.duration_seconds:.1f}s, {response.video_content.resolution}, {response.video_content.format}\\par")
            
            rtf_content.append(r"\par")
        
        # Citations
        if response.knowledge_citations:
            rtf_content.append(r"\fs20\b Citations\b0\fs18\par")
            
            for i, citation in enumerate(response.knowledge_citations, 1):
                title_escaped = citation.source_title.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}')
                location_escaped = citation.location_reference.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}')
                
                citation_text = f"{i}. {title_escaped} ({citation.source_type.value.title()}) - {location_escaped}\\par"
                rtf_content.append(citation_text)
        
        # RTF footer
        rtf_content.append("}")
        
        return "\n".join(rtf_content).encode('utf-8')
    
    def _export_to_pptx(self, response: MultimediaResponse) -> bytes:
        """Export to PowerPoint presentation with slides."""
        logger.debug("Exporting to PPTX format")
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Multimodal Librarian Export"
        subtitle.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Content slide
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Content"
        content.text = response.text_content
        
        # Visualizations slides
        if response.visualizations:
            for i, viz in enumerate(response.visualizations, 1):
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
                
                # Add title
                title_shape = slide.shapes.add_textbox(PPTXInches(0.5), PPTXInches(0.5), PPTXInches(9), PPTXInches(1))
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = f"Visualization {i}: {viz.viz_type.title()}"
                title_para.font.size = PPTXPt(24)
                title_para.font.bold = True
                
                # Add image if available
                if viz.content_data:
                    try:
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                            tmp_file.write(viz.content_data)
                            tmp_file.flush()
                            
                            # Add image to slide
                            slide.shapes.add_picture(tmp_file.name, PPTXInches(1), PPTXInches(2), PPTXInches(8), PPTXInches(5))
                            
                            os.unlink(tmp_file.name)
                    except Exception as e:
                        logger.warning(f"Could not embed visualization {i}: {e}")
                        # Add placeholder text
                        placeholder_shape = slide.shapes.add_textbox(PPTXInches(1), PPTXInches(3), PPTXInches(8), PPTXInches(2))
                        placeholder_frame = placeholder_shape.text_frame
                        placeholder_frame.text = f"[{viz.viz_type.title()} Visualization]"
                
                # Add caption
                if viz.caption:
                    caption_shape = slide.shapes.add_textbox(PPTXInches(1), PPTXInches(7.5), PPTXInches(8), PPTXInches(1))
                    caption_frame = caption_shape.text_frame
                    caption_para = caption_frame.paragraphs[0]
                    caption_para.text = viz.caption
                    caption_para.font.size = PPTXPt(14)
                    caption_para.font.italic = True
                    caption_para.alignment = PP_ALIGN.CENTER
        
        # Citations slide
        if response.knowledge_citations:
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Citations"
            
            citations_text = ""
            for i, citation in enumerate(response.knowledge_citations, 1):
                citations_text += f"{i}. {citation.source_title}\n"
                citations_text += f"   Source: {citation.source_type.value.title()}\n"
                citations_text += f"   Location: {citation.location_reference}\n"
                citations_text += f"   Relevance: {citation.relevance_score:.2f}\n\n"
            
            content.text = citations_text
        
        # Save to bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer.getvalue()
    
    def _export_to_xlsx(self, response: MultimediaResponse) -> bytes:
        """Export to Excel spreadsheet with structured data."""
        logger.debug("Exporting to XLSX format")
        
        wb = Workbook()
        
        # Remove default sheet and create custom sheets
        wb.remove(wb.active)
        
        # Content sheet
        content_ws = wb.create_sheet("Content")
        content_ws['A1'] = "Multimodal Librarian Export"
        content_ws['A1'].font = Font(size=16, bold=True)
        content_ws['A2'] = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        content_ws['A4'] = "Main Content:"
        content_ws['A4'].font = Font(bold=True)
        content_ws['A5'] = response.text_content
        content_ws['A5'].alignment = Alignment(wrap_text=True, vertical='top')
        
        # Adjust column width
        content_ws.column_dimensions['A'].width = 80
        
        # Visualizations sheet
        if response.visualizations:
            viz_ws = wb.create_sheet("Visualizations")
            viz_ws['A1'] = "Visualizations"
            viz_ws['A1'].font = Font(size=14, bold=True)
            
            headers = ['ID', 'Type', 'Caption', 'Alt Text']
            for col, header in enumerate(headers, 1):
                cell = viz_ws.cell(row=3, column=col, value=header)
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
            
            for row, viz in enumerate(response.visualizations, 4):
                viz_ws.cell(row=row, column=1, value=viz.viz_id)
                viz_ws.cell(row=row, column=2, value=viz.viz_type)
                viz_ws.cell(row=row, column=3, value=viz.caption)
                viz_ws.cell(row=row, column=4, value=viz.alt_text)
            
            # Adjust column widths
            for col in range(1, 5):
                viz_ws.column_dimensions[chr(64 + col)].width = 20
        
        # Citations sheet
        if response.knowledge_citations:
            citations_ws = wb.create_sheet("Citations")
            citations_ws['A1'] = "Citations"
            citations_ws['A1'].font = Font(size=14, bold=True)
            
            headers = ['#', 'Title', 'Source Type', 'Location', 'Relevance Score']
            for col, header in enumerate(headers, 1):
                cell = citations_ws.cell(row=3, column=col, value=header)
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
            
            for row, citation in enumerate(response.knowledge_citations, 4):
                citations_ws.cell(row=row, column=1, value=row - 3)
                citations_ws.cell(row=row, column=2, value=citation.source_title)
                citations_ws.cell(row=row, column=3, value=citation.source_type.value.title())
                citations_ws.cell(row=row, column=4, value=citation.location_reference)
                citations_ws.cell(row=row, column=5, value=citation.relevance_score)
            
            # Adjust column widths
            citations_ws.column_dimensions['A'].width = 5
            citations_ws.column_dimensions['B'].width = 40
            citations_ws.column_dimensions['C'].width = 15
            citations_ws.column_dimensions['D'].width = 20
            citations_ws.column_dimensions['E'].width = 15
        
        # Multimedia sheet
        if response.audio_content or response.video_content:
            media_ws = wb.create_sheet("Multimedia")
            media_ws['A1'] = "Multimedia Content"
            media_ws['A1'].font = Font(size=14, bold=True)
            
            row = 3
            if response.audio_content:
                media_ws[f'A{row}'] = "Audio Content:"
                media_ws[f'A{row}'].font = Font(bold=True)
                row += 1
                media_ws[f'A{row}'] = f"Duration: {response.audio_content.duration_seconds:.1f} seconds"
                row += 1
                media_ws[f'A{row}'] = f"Format: {response.audio_content.format}"
                row += 2
            
            if response.video_content:
                media_ws[f'A{row}'] = "Video Content:"
                media_ws[f'A{row}'].font = Font(bold=True)
                row += 1
                media_ws[f'A{row}'] = f"Duration: {response.video_content.duration_seconds:.1f} seconds"
                row += 1
                media_ws[f'A{row}'] = f"Resolution: {response.video_content.resolution}"
                row += 1
                media_ws[f'A{row}'] = f"Format: {response.video_content.format}"
            
            media_ws.column_dimensions['A'].width = 50
        
        # Save to bytes
        xlsx_buffer = io.BytesIO()
        wb.save(xlsx_buffer)
        xlsx_buffer.seek(0)
        return xlsx_buffer.getvalue()
    
    def preserve_formatting(self, content: Any, format_type: str) -> Any:
        """
        Preserve visual integrity across formats.
        
        Args:
            content: Content to format
            format_type: Target format
            
        Returns:
            Formatted content appropriate for the target format
        """
        # This method can be extended to handle format-specific formatting preservation
        # For now, it returns the content as-is since formatting is handled in individual export methods
        return content
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported export formats."""
        return list(self.supported_formats.keys())
    
    def validate_response(self, response: MultimediaResponse) -> bool:
        """
        Validate that the response can be exported.
        
        Args:
            response: MultimediaResponse to validate
            
        Returns:
            bool: True if response is valid for export
        """
        if not response or not response.text_content:
            return False
        
        return response.validate()
    
    def _export_to_html(self, response: MultimediaResponse) -> bytes:
        """Export to HTML for web-compatible output."""
        logger.debug("Exporting to HTML format")
        
        import base64
        
        html_parts = []
        
        # HTML header with CSS styling
        html_parts.append("""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Multimodal Librarian Export</title>
    <style>
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, sans-serif;
            line-height: 1.6;
            max-width: 900px;
            margin: 0 auto;
            padding: 20px;
            color: #333;
            background-color: #f9f9f9;
        }
        h1 {
            color: #2c3e50;
            border-bottom: 2px solid #3498db;
            padding-bottom: 10px;
        }
        h2 {
            color: #34495e;
            margin-top: 30px;
        }
        h3 {
            color: #7f8c8d;
        }
        .metadata {
            color: #7f8c8d;
            font-size: 0.9em;
            margin-bottom: 20px;
        }
        .content {
            background-color: white;
            padding: 20px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            margin-bottom: 20px;
        }
        .visualization {
            text-align: center;
            margin: 20px 0;
            padding: 15px;
            background-color: white;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        }
        .visualization img {
            max-width: 100%;
            height: auto;
            border-radius: 4px;
        }
        .caption {
            font-style: italic;
            color: #7f8c8d;
            margin-top: 10px;
        }
        .citation {
            background-color: white;
            padding: 15px;
            margin: 10px 0;
            border-left: 4px solid #3498db;
            border-radius: 0 8px 8px 0;
        }
        .citation-title {
            font-weight: bold;
            color: #2c3e50;
        }
        .citation-details {
            font-size: 0.9em;
            color: #7f8c8d;
        }
        .multimedia-info {
            background-color: #ecf0f1;
            padding: 15px;
            border-radius: 8px;
            margin: 10px 0;
        }
        .relevance-score {
            display: inline-block;
            background-color: #3498db;
            color: white;
            padding: 2px 8px;
            border-radius: 4px;
            font-size: 0.8em;
        }
    </style>
</head>
<body>
""")
        
        # Title and metadata
        html_parts.append(f"""
    <h1>Multimodal Librarian Export</h1>
    <div class="metadata">Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</div>
""")
        
        # Main content
        html_parts.append("""
    <h2>Content</h2>
    <div class="content">
""")
        # Escape HTML special characters and preserve line breaks
        content_escaped = response.text_content.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        content_escaped = content_escaped.replace('\n', '<br>\n')
        html_parts.append(f"        <p>{content_escaped}</p>\n")
        html_parts.append("    </div>\n")
        
        # Visualizations
        if response.visualizations:
            html_parts.append("    <h2>Visualizations</h2>\n")
            
            for i, viz in enumerate(response.visualizations, 1):
                html_parts.append(f"""
    <div class="visualization">
        <h3>{i}. {viz.viz_type.title()}</h3>
""")
                
                # Embed image as base64 if available
                if viz.content_data:
                    img_base64 = base64.b64encode(viz.content_data).decode('utf-8')
                    html_parts.append(f'        <img src="data:image/png;base64,{img_base64}" alt="{viz.alt_text or viz.viz_type}">\n')
                else:
                    html_parts.append(f'        <p>[{viz.viz_type.title()} Visualization]</p>\n')
                
                if viz.caption:
                    caption_escaped = viz.caption.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    html_parts.append(f'        <div class="caption">{caption_escaped}</div>\n')
                
                html_parts.append("    </div>\n")
        
        # Multimedia content info
        if response.audio_content or response.video_content:
            html_parts.append("    <h2>Multimedia Content</h2>\n")
            
            if response.audio_content:
                html_parts.append(f"""
    <div class="multimedia-info">
        <h3>Audio</h3>
        <p>Duration: {response.audio_content.duration_seconds:.1f} seconds</p>
        <p>Format: {response.audio_content.format}</p>
    </div>
""")
            
            if response.video_content:
                html_parts.append(f"""
    <div class="multimedia-info">
        <h3>Video</h3>
        <p>Duration: {response.video_content.duration_seconds:.1f} seconds</p>
        <p>Resolution: {response.video_content.resolution}</p>
        <p>Format: {response.video_content.format}</p>
    </div>
""")
        
        # Citations
        if response.knowledge_citations:
            html_parts.append("    <h2>Citations</h2>\n")
            
            for i, citation in enumerate(response.knowledge_citations, 1):
                title_escaped = citation.source_title.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                location_escaped = citation.location_reference.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                
                html_parts.append(f"""
    <div class="citation">
        <div class="citation-title">{i}. {title_escaped}</div>
        <div class="citation-details">
            Source: {citation.source_type.value.title()}<br>
            Location: {location_escaped}<br>
            <span class="relevance-score">Relevance: {citation.relevance_score:.2f}</span>
        </div>
    </div>
""")
        
        # HTML footer
        html_parts.append("""
</body>
</html>
""")
        
        return "".join(html_parts).encode('utf-8')